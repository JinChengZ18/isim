#!/usr/bin/env python3
"""Compose the Chapter-3 circuit-section figures from clean panels into a
PowerPoint deck, add panel letters (a)(b)(c)(d), and export numbered PNGs.

This mirrors the smtj_pbnn_sim pipeline: generator scripts emit only clean,
letter-free panels (here eda/figs_make.py -> eda/figs_raw/); this script is the
single place where panel letters and figure numbering live. It builds one slide
per figure in article/ppt/Chapter03_local.pptx, drops an (x) text box at each
panel's top-left, then renders the deck via LibreOffice and exports each slide,
auto-cropped, to article/figs/Chapter03_local_NN.png.

Run:  python eda/build_ppt_figs.py            # (after eda/figs_make.py)
Deps: python-pptx, PyMuPDF (fitz), Pillow, LibreOffice at SOFFICE.
"""
from __future__ import annotations

import subprocess
from pathlib import Path

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
RAW = HERE / "figs_raw"
PPT = ROOT / "article" / "ppt"
FIGS = ROOT / "article" / "figs"
DECK = PPT / "Chapter03_local.pptx"

SOFFICE = r"C:\Program Files\LibreOffice\program\soffice.exe"

# slide + layout geometry (inches)
SLIDE_W, SLIDE_H = 10.0, 13.0
MARGIN, GAP, TOP = 0.25, 0.16, 0.42
LETTER_DX, LETTER_DY = 0.04, 0.02   # letter offset from panel top-left
EXPORT_DPI = 419
CROP_PAD = 8

# one slide per figure; panels named <key>.png in figs_raw/, letters assigned
# left-to-right, top-to-bottom in caption reading order.
FIGURES = [
    {"num": 9,  "rows": [["chain_schematic"],
                         ["chain_transfer", "chain_waveform"],
                         ["chain_corners"]]},
    {"num": 10, "rows": [["abl_bits", "abl_span", "abl_reset"],
                         ["abl_traj", "abl_readflip"]]},
    {"num": 11, "rows": [["ir_profile", "ir_impact"],
                         ["energy_stack", "hw_projection"]]},
]


def aspect(key):
    """height/width of the clean panel PNG."""
    with Image.open(RAW / f"{key}.png") as im:
        w, h = im.size
    return h / w


def letter_box(slide, lx, ly, ch):
    tb = slide.shapes.add_textbox(Inches(lx), Inches(ly), Inches(0.5),
                                  Inches(0.30))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    r = tf.paragraphs[0].add_run()
    r.text = f"({ch})"
    r.font.name = "Arial"
    r.font.bold = True
    r.font.size = Pt(15)
    r.font.color.rgb = RGBColor(0x20, 0x20, 0x20)


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    return slide


def layout_figure(slide, rows):
    """Place panels row by row, equal widths within a row, tops aligned; drop a
    letter at each panel's top-left. Returns the used content bbox height."""
    tw = SLIDE_W - 2 * MARGIN
    letters = "abcdefgh"
    li = 0
    y = TOP
    for row in rows:
        n = len(row)
        wp = (tw - (n - 1) * GAP) / n
        heights = [wp * aspect(k) for k in row]
        row_h = max(heights)
        x = MARGIN
        for k, h in zip(row, heights):
            slide.shapes.add_picture(str(RAW / f"{k}.png"), Inches(x),
                                     Inches(y), width=Inches(wp),
                                     height=Inches(h))
            letter_box(slide, x + LETTER_DX, y + LETTER_DY, letters[li])
            li += 1
            x += wp + GAP
        y += row_h + GAP
    return y


def crop_content(im):
    from PIL import ImageChops
    bg = Image.new("RGB", im.size, (255, 255, 255))
    diff = ImageChops.difference(im.convert("RGB"), bg)
    bbox = diff.getbbox()
    if not bbox:
        return im
    l, t, r, b = bbox
    l = max(0, l - CROP_PAD); t = max(0, t - CROP_PAD)
    r = min(im.width, r + CROP_PAD); b = min(im.height, b + CROP_PAD)
    return im.crop((l, t, r, b))


def build():
    PPT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    for fig in FIGURES:
        slide = blank_slide(prs)
        layout_figure(slide, fig["rows"])
    prs.save(str(DECK))
    print(f"deck -> {DECK}")


def export():
    subprocess.run([SOFFICE, "--headless", "--convert-to", "pdf",
                    "--outdir", str(PPT), str(DECK)], check=True,
                   capture_output=True)
    pdf = fitz.open(str(PPT / "Chapter03_local.pdf"))
    for i, fig in enumerate(FIGURES):
        page = pdf.load_page(i)
        pix = page.get_pixmap(dpi=EXPORT_DPI)
        im = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
        im = crop_content(im)
        out = FIGS / f"Chapter03_local_{fig['num']:02d}.png"
        im.save(out)
        print(f"figure -> {out}  ({im.width}x{im.height})")
    pdf.close()


if __name__ == "__main__":
    build()
    export()
